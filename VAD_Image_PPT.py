from selenium import webdriver
from selenium.webdriver import ActionChains
from selenium.webdriver.common.by import By
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.chrome.service import Service
from pptx import Presentation
from pptx.util import Inches
from time import sleep
from time import gmtime, strftime
import os

cwd = os.getcwd()
service = Service(cwd + "\chromedriver.exe")
options = webdriver.ChromeOptions()
options = Options()
options.add_argument("start-maximized")
options.headless = True
driver = webdriver.Chrome(service=service, options=options)
driver.set_window_size(1920, 1080)

#INPUTS

Event = 'Straßenfest'
ISO_Name = 'DEU'
Coord ='52.141153,9.93461'.split(",")
Zoom = '16'
Lat = Coord[0]
Lon = Coord[1]


Zoom_TT_int = int(Zoom)-1.25
Zoom_TT = str(Zoom_TT_int)

filename_mydrive = 'TomTom_' + (strftime("%Y-%m-%d %H%M", gmtime()))

def mydrive (Lat=Lat,Lon=Lon,Zoom_TT=Zoom_TT):
    url = 'https://mydrive.tomtom.com/en_gb/#mode=viewport+viewport=' + Lat +',' + Lon + ',' + Zoom_TT + ',0,-0+ver=3'
    driver.get(url)
    sleep(20)
    #driver.implicitly_wait(100)
    driver.get_screenshot_as_file(filename_mydrive + '.png')   

filename_plan = 'Plan_' + (strftime("%Y-%m-%d %H%M", gmtime()))

def plan (Lat=Lat,Lon=Lon,Zoom=Zoom):
    url = 'https://plan.tomtom.com/en/?p=' + Lat +',' + Lon + ',' + Zoom + 'z'
    driver.get(url)
    sleep(10)
    #driver.implicitly_wait(100)
    driver.get_screenshot_as_file(filename_plan + '.png') 

filename_googlemaps = 'GoogleMaps_' + (strftime("%Y-%m-%d %H%M", gmtime()))

def googlemaps (Lat=Lat,Lon=Lon,Zoom=Zoom):
    url = 'https://www.google.com/maps/@' + Lat + ',' + Lon + ',' + Zoom + 'z/data=!5m1!1e1'
    driver.get(url)
    sleep(10)
    driver.find_element(by=By.XPATH, value='//span[text()="Accept all"]').click()
    sleep(10)
    driver.implicitly_wait(100)
    driver.get_screenshot_as_file(filename_googlemaps + '.png')      

filename_waze = 'Waze_' + (strftime("%Y-%m-%d %H%M", gmtime()))

def waze (Lat=Lat,Lon=Lon,Zoom=Zoom):
    url = 'https://embed.waze.com/iframe?' + 'zoom=' + Zoom + '&lat=' + Lat + '&lon=' + Lon + '&ct=livemap'
    driver.get(url)
    sleep(10)
    #driver.implicitly_wait(100)
    driver.get_screenshot_as_file(filename_waze + '.png')     

filename_here = 'Here_' + (strftime("%Y-%m-%d %H%M", gmtime()))

def here (Lat=Lat,Lon=Lon,Zoom=Zoom):
    url = 'https://wego.here.com/traffic/explore?map=' + Lat + ',' + Lon + ',' + Zoom + ',traffic'
    driver.get(url)
    sleep(10)
    driver.find_element(by=By.XPATH, value='/html/body/div[2]/div/div[2]/button').click()
    sleep(10)
    driver.implicitly_wait(100)
    driver.get_screenshot_as_file(filename_here + '.png')   

mydrive()
#plan()
googlemaps()
waze()
here()

#Create blank PowerPoint file
prs = Presentation()

#Define Images
TomTom_png = filename_mydrive + '.png'
Here_png = filename_here + '.png'
Waze_png = filename_waze + '.png'
GoogleMaps_png = filename_googlemaps + '.png'

#Event Name Input
Event_Name = Event

#Add a slide
slide = prs.slides.add_slide(prs.slide_layouts[5])
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

#Add slide title
slide.shapes[0].text = Event_Name


#Add slide images
#Location, Left, Top, Height, Width
TomTom = slide.shapes.add_picture(TomTom_png,Inches(0.2),Inches(1.89),height=Inches(4.78),width=Inches(8.49))
TomTom_Logo = slide.shapes.add_picture('Logo_TomTom.jpg',Inches(0.2),Inches(1.89))
Waze = slide.shapes.add_picture(Waze_png,Inches(8.9),Inches(0.11),height=Inches(2.39),width=Inches(4.24))
Waze_Logo = slide.shapes.add_picture('Logo_Waze.jpg',Inches(8.9),Inches(0.11))
Here=slide.shapes.add_picture(Here_png,Inches(8.9),Inches(2.61),height=Inches(2.39),width=Inches(4.24))
Here_Logo = slide.shapes.add_picture('Logo_Here.jpg',Inches(8.9),Inches(2.61))
GoogleMaps = slide.shapes.add_picture(GoogleMaps_png,Inches(8.9),Inches(5.11),height=Inches(2.39),width=Inches(4.24))
GoogleMaps_Logo = slide.shapes.add_picture('Logo_GoogleMaps.jpg',Inches(8.9),Inches(5.11))

prs.save(ISO_Name+"_"+Event_Name+".pptx")
